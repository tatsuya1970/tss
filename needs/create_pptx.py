from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── TSS テレビ新広島 公式ブランドカラー ──────────────────
TSS_GREEN      = RGBColor(0x00, 0x55, 0x3F)  # ダークグリーン（サイトのメインカラー）
TSS_GREEN_LIGHT= RGBColor(0x00, 0x6A, 0x4F)  # グリーン明るめ
TSS_RED        = RGBColor(0xD5, 0x2E, 0x31)  # ブランドレッド
TSS_BLUE       = RGBColor(0x1D, 0x78, 0xAB)  # ブルー
TSS_GOLD       = RGBColor(0xEF, 0xBA, 0x31)  # ゴールド
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK           = RGBColor(0x33, 0x33, 0x33)
GRAY           = RGBColor(0x70, 0x70, 0x70)
LIGHT_BG       = RGBColor(0xF5, 0xF5, 0xF5)
BORDER         = RGBColor(0xE4, 0xE4, 0xE4)
GREEN_BG       = RGBColor(0xE8, 0xF4, 0xF0)  # 薄グリーン背景

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=0.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    # 上部アクセントライン（細い）
    rect(slide, 0, 0, 13.33, 0.08, fill=TSS_RED)
    # ヘッダー背景
    rect(slide, 0, 0.08, 13.33, 1.12, fill=TSS_GREEN)
    txt(slide, title, 0.45, 0.15, 10, 0.65,
        size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.78, 11, 0.38,
            size=13, color=RGBColor(0xA8, 0xD5, 0xC2), italic=True)


def page_num(slide, n):
    txt(slide, f"{n}", 12.85, 7.1, 0.45, 0.35,
        size=11, color=GRAY, align=PP_ALIGN.RIGHT)


def badge(slide, label, l, t, w=1.1, h=0.38, color=TSS_GREEN):
    rect(slide, l, t, w, h, fill=color)
    txt(slide, label, l, t, w, h,
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# スライド 1: 表紙
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# 上半分: TSS グリーン
rect(slide, 0, 0, 13.33, 4.5, fill=TSS_GREEN)
# 下半分: 白
rect(slide, 0, 4.5, 13.33, 3.0, fill=WHITE)
# 赤いアクセントライン
rect(slide, 0, 4.4, 13.33, 0.12, fill=TSS_RED)
# ゴールドの細線
rect(slide, 0, 4.52, 13.33, 0.04, fill=TSS_GOLD)

txt(slide, "株式会社テレビ新広島 御中", 1.0, 0.6, 11.33, 0.6,
    size=16, color=RGBColor(0xA8, 0xD5, 0xC2), align=PP_ALIGN.CENTER)
txt(slide, "AI活用による地域情報収集・SNS分析", 0.8, 1.3, 11.73, 0.9,
    size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "ソリューションのご提案", 0.8, 2.15, 11.73, 0.8,
    size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ゴールドの装飾ライン
rect(slide, 3.5, 3.2, 6.33, 0.05, fill=TSS_GOLD)

txt(slide, "2026年3月26日", 0.8, 5.1, 11.73, 0.6,
    size=15, color=GRAY, align=PP_ALIGN.CENTER)
txt(slide, "Presented by TSS AIソリューション提案チーム", 0.8, 5.7, 11.73, 0.5,
    size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════
# スライド 2: 課題認識
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
rect(slide, 0, 0, 13.33, 0.06, fill=TSS_RED)
header(slide, "貴社の課題認識", "現状の情報収集に潜む3つの壁")
page_num(slide, 2)

issues = [
    ("①", "人力での情報収集に限界",
     "広島県23市町の最新情報を\nリサーチャーが毎回手動確認。\n見落としや確認漏れが発生。"),
    ("②", "SNSデータを\n報道に活かせていない",
     "地域の声・リアルタイム情報が\n体系的に把握できておらず、\n活用の手が止まっている。"),
    ("③", "情報の一元管理ツールが\n存在しない",
     "情報が分散しており、\n優先度の判断やチームでの\n迅速な共有が困難。"),
]

for i, (num, title, desc) in enumerate(issues):
    x = 0.45 + i * 4.25
    # カード外枠
    rect(slide, x, 1.45, 3.95, 5.5, fill=WHITE, line_color=BORDER, line_w=1)
    # 上部グリーンヘッダー
    rect(slide, x, 1.45, 3.95, 1.0, fill=TSS_GREEN)
    txt(slide, num, x, 1.5, 3.95, 0.9,
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # タイトル帯
    rect(slide, x, 2.45, 3.95, 0.05, fill=TSS_RED)
    txt(slide, title, x + 0.15, 2.55, 3.65, 1.1,
        size=16, bold=True, color=TSS_GREEN)
    txt(slide, desc, x + 0.15, 3.75, 3.65, 2.2,
        size=13, color=GRAY)


# ═══════════════════════════════════════════════════════════
# スライド 3: ソリューション全体像
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
rect(slide, 0, 0, 13.33, 0.06, fill=TSS_RED)
header(slide, "ソリューション全体像", "2つのAIシステムで3つの課題を解決します")
page_num(slide, 3)

# 左カード
rect(slide, 0.4, 1.4, 5.85, 5.5, fill=GREEN_BG, line_color=TSS_GREEN, line_w=1.5)
rect(slide, 0.4, 1.4, 5.85, 0.6, fill=TSS_GREEN)
txt(slide, "① 広島市町情報ダッシュボード", 0.55, 1.43, 5.55, 0.55,
    size=14, bold=True, color=WHITE)
items1 = [
    "全23市町の公式サイトを自動巡回",
    "AIが記事を要約・分類・スコアリング",
    "カテゴリ・期間・市町でフィルター表示",
    "過去記事の蓄積・時系列分析",
]
for j, item in enumerate(items1):
    txt(slide, "─  " + item, 0.6, 2.2 + j * 0.72, 5.45, 0.65,
        size=13, color=DARK)

rect(slide, 0.55, 5.55, 1.55, 0.42, fill=TSS_RED)
txt(slide, "✅  稼働中・デモ可能", 0.55, 5.55, 3.5, 0.42,
    size=12, bold=True, color=WHITE)

# 中央の「+」または統合メモ
rect(slide, 6.4, 3.1, 0.55, 0.55, fill=TSS_GOLD)
txt(slide, "+", 6.4, 3.1, 0.55, 0.55,
    size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "統合も\n可能", 6.25, 3.75, 0.85, 0.7,
    size=10, bold=True, color=TSS_RED, align=PP_ALIGN.CENTER)

# 右カード
rect(slide, 7.1, 1.4, 5.85, 5.5, fill=RGBColor(0xFD, 0xF2, 0xF2),
     line_color=TSS_RED, line_w=1.5)
rect(slide, 7.1, 1.4, 5.85, 0.6, fill=TSS_RED)
txt(slide, "② SNSトレンド可視化ツール", 7.25, 1.43, 5.55, 0.55,
    size=14, bold=True, color=WHITE)
items2 = [
    "X（旧Twitter）のリアルタイム投稿を収集",
    "地図上にヒートマップで可視化",
    "AIが他県地名の誤検知を自動除外",
    "Instagram連携の実装準備済み",
]
for j, item in enumerate(items2):
    txt(slide, "─  " + item, 7.3, 2.2 + j * 0.72, 5.5, 0.65,
        size=13, color=DARK)

rect(slide, 7.25, 5.55, 1.55, 0.42, fill=TSS_RED)
txt(slide, "✅  稼働中・デモ可能", 7.25, 5.55, 3.5, 0.42,
    size=12, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════
# スライド 4: ダッシュボード詳細
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
rect(slide, 0, 0, 13.33, 0.06, fill=TSS_RED)
header(slide, "① 広島市町情報ダッシュボード", "全23市町の一次情報をAIが24時間自動収集・要約")
page_num(slide, 4)

# 左：対応市町
rect(slide, 0.4, 1.4, 5.85, 2.7, fill=GREEN_BG, line_color=TSS_GREEN, line_w=1)
rect(slide, 0.4, 1.4, 5.85, 0.48, fill=TSS_GREEN)
txt(slide, "対応市町（全23市町）", 0.55, 1.43, 5.55, 0.43,
    size=13, bold=True, color=WHITE)
cities = ("広島市・福山市・呉市・廿日市市・尾道市・東広島市・竹原市・三原市・"
          "府中市・三次市・庄原市・大竹市・安芸高田市・江田島市・府中町・海田町・"
          "熊野町・坂町・安芸太田町・北広島町・大崎上島町・世羅町・神石高原町")
txt(slide, cities, 0.55, 2.0, 5.65, 2.0, size=11, color=DARK)

# 左：AIカテゴリ
rect(slide, 0.4, 4.25, 5.85, 2.65, fill=GREEN_BG, line_color=TSS_GREEN, line_w=1)
rect(slide, 0.4, 4.25, 5.85, 0.48, fill=TSS_GREEN)
txt(slide, "AIによる自動分類カテゴリ", 0.55, 4.28, 5.55, 0.43,
    size=13, bold=True, color=WHITE)
cats = [("防災・安全", TSS_RED), ("福祉・医療", TSS_BLUE),
        ("経済・産業", TSS_GREEN), ("教育・文化", TSS_GOLD),
        ("インフラ・都市", GRAY), ("その他", GRAY)]
for ci, (cname, ccol) in enumerate(cats):
    bx = 0.55 + (ci % 3) * 1.85
    by = 4.88 + (ci // 3) * 0.6
    rect(slide, bx, by, 1.65, 0.45, fill=ccol)
    txt(slide, cname, bx, by, 1.65, 0.45,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 右：機能
rect(slide, 6.65, 1.4, 6.3, 2.7, fill=WHITE, line_color=BORDER, line_w=1)
rect(slide, 6.65, 1.4, 6.3, 0.48, fill=TSS_GREEN)
txt(slide, "主な機能", 6.8, 1.43, 6.0, 0.43,
    size=13, bold=True, color=WHITE)
funcs = [
    "新着情報を自動スクレイピング（重複排除）",
    "タイトルベースのAI要約＋重要度スコア（1〜5）",
    "市町・カテゴリ・日付範囲でのフィルター",
    "過去記事の蓄積・トレンド時系列分析",
]
for j, f in enumerate(funcs):
    txt(slide, "✓  " + f, 6.8, 2.02 + j * 0.5, 6.0, 0.46,
        size=12, color=DARK)

# 右：ロードマップ
rect(slide, 6.65, 4.25, 6.3, 2.65, fill=WHITE, line_color=BORDER, line_w=1)
rect(slide, 6.65, 4.25, 6.3, 0.48, fill=TSS_GOLD)
txt(slide, "今後の拡張ロードマップ", 6.8, 4.28, 6.0, 0.43,
    size=13, bold=True, color=WHITE)
roadmap = [
    ("Phase 2", "広報資料（PDF等）の本文取得と詳細要約"),
    ("Phase 2", "オープンデータポータルとの連携"),
    ("Phase 3", "記事本文全文を読んだ高精度な要約"),
]
for j, (ph, desc) in enumerate(roadmap):
    y = 4.88 + j * 0.63
    badge(slide, ph, 6.8, y, w=0.9, h=0.4, color=TSS_GOLD)
    txt(slide, desc, 7.82, y, 4.9, 0.4, size=12, color=DARK)


# ═══════════════════════════════════════════════════════════
# スライド 5: SNS詳細
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
rect(slide, 0, 0, 13.33, 0.06, fill=TSS_RED)
header(slide, "② SNSトレンド可視化ツール", "現場のリアルタイムな声を地図上で把握")
page_num(slide, 5)

# Xの価値帯
rect(slide, 0.4, 1.4, 12.53, 1.45, fill=GREEN_BG, line_color=TSS_GREEN, line_w=1)
rect(slide, 0.4, 1.4, 0.06, 1.45, fill=TSS_GREEN)
txt(slide, "X（旧Twitter）が持つ唯一の価値 ── リアルタイム性",
    0.6, 1.48, 11.8, 0.5, size=15, bold=True, color=TSS_GREEN)
txt(slide,
    "情報量は少なくても「公式発表より早く現場の声が上がる」という点は他SNSで代替不可。"
    "イベント当日の混雑・速報・地域の盛り上がりを、いち早くキャッチできる。",
    0.6, 2.0, 12.0, 0.72, size=13, color=DARK)

# 左：機能
rect(slide, 0.4, 3.05, 5.85, 3.85, fill=WHITE, line_color=BORDER, line_w=1)
rect(slide, 0.4, 3.05, 5.85, 0.48, fill=TSS_RED)
txt(slide, "主な機能", 0.55, 3.08, 5.55, 0.43,
    size=13, bold=True, color=WHITE)
sns_funcs = [
    "地図ヒートマップ表示（熱量を直感的に把握）",
    "投稿数・いいね・RT数から盛り上がりスコア算出",
    "GPT-4oによる他県同名地名の誤検知除外",
    "イベント・グルメ・交通トラブル等を自動分類",
    "元投稿URLをワンクリックで確認",
]
for j, f in enumerate(sns_funcs):
    txt(slide, "✓  " + f, 0.6, 3.65 + j * 0.58, 5.5, 0.52,
        size=12, color=DARK)

# 右：SNS拡張
rect(slide, 6.65, 3.05, 6.3, 3.85, fill=WHITE, line_color=BORDER, line_w=1)
rect(slide, 6.65, 3.05, 6.3, 0.48, fill=TSS_RED)
txt(slide, "SNS拡張ロードマップ", 6.8, 3.08, 6.0, 0.43,
    size=13, bold=True, color=WHITE)

platforms = [
    (TSS_GREEN,  "稼働中",   "X（旧Twitter）",   "リアルタイム・速報性"),
    (TSS_GOLD,   "準備済み", "Instagram",         "グルメ・観光・イベントの盛り上がり"),
    (GRAY,       "今後対応", "TikTok",            "若年層の地域情報"),
    (GRAY,       "今後対応", "Googleクチコミ",     "飲食・観光スポットの声"),
]
for j, (col, status, name, desc) in enumerate(platforms):
    y = 3.65 + j * 0.77
    badge(slide, status, 6.8, y, w=0.95, h=0.42, color=col)
    txt(slide, name,  7.87, y,        1.9, 0.42, size=13, bold=True, color=DARK)
    txt(slide, desc,  9.88, y + 0.02, 2.9, 0.38, size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════
# スライド 6: 統合オプション
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
rect(slide, 0, 0, 13.33, 0.06, fill=TSS_RED)
header(slide, "統合オプション", "2システムを1つの統合ダッシュボードにまとめることも可能です")
page_num(slide, 6)

# 左：分離
rect(slide, 0.4, 1.4, 5.85, 5.0, fill=WHITE, line_color=TSS_GREEN, line_w=1.5)
rect(slide, 0.4, 1.4, 5.85, 0.55, fill=TSS_GREEN)
txt(slide, "分離構成（現行）", 0.55, 1.43, 5.55, 0.5,
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
sep = [
    "市町情報ダッシュボード（独立運用）",
    "SNSトレンドツール（独立運用）",
    "それぞれ単独での導入・運用が可能",
    "段階的な導入・予算分割に適している",
    "既存システムへの影響が少ない",
]
for j, s in enumerate(sep):
    txt(slide, "─  " + s, 0.6, 2.15 + j * 0.72, 5.45, 0.65,
        size=13, color=DARK)

# 中央矢印
txt(slide, "または", 6.3, 3.65, 0.73, 0.5,
    size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# 右：統合
rect(slide, 7.1, 1.4, 5.85, 5.0, fill=RGBColor(0xFD, 0xF2, 0xF2),
     line_color=TSS_RED, line_w=2)
rect(slide, 7.1, 1.4, 5.85, 0.55, fill=TSS_RED)
txt(slide, "統合構成（検討可能）", 7.25, 1.43, 5.55, 0.5,
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
uni = [
    "公式新着情報＋SNSトレンドを同一画面に集約",
    "広島県版ニュース統合プラットフォーム",
    "リサーチャーが1画面で全情報を確認",
    "将来的な機能追加・拡張に対応しやすい",
    "AIアシスト機能の横断活用が可能",
]
for j, u in enumerate(uni):
    txt(slide, "─  " + u, 7.3, 2.15 + j * 0.72, 5.5, 0.65,
        size=13, color=DARK)

# 下部メッセージ
rect(slide, 0.4, 6.6, 12.53, 0.7, fill=GREEN_BG)
txt(slide,
    "導入規模・運用体制・ご予算に合わせて最適な構成をご提案します。ヒアリングの場でご要望をお聞かせください。",
    0.6, 6.68, 12.1, 0.52, size=13, color=TSS_GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# スライド 7: 導入効果 Before/After
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
rect(slide, 0, 0, 13.33, 0.06, fill=TSS_RED)
header(slide, "導入効果", "Before / After")
page_num(slide, 7)

# テーブルヘッダー
col_x = [0.4, 3.3, 7.5]
col_w = [2.75, 4.05, 5.5]
headers_txt = ["課題", "現状（Before）", "導入後（After）"]
hcols = [TSS_GREEN, GRAY, TSS_GREEN]
for i, (h, x, w, hc) in enumerate(zip(headers_txt, col_x, col_w, hcols)):
    rect(slide, x, 1.4, w, 0.5, fill=hc)
    txt(slide, h, x + 0.1, 1.43, w - 0.2, 0.45,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("市町情報\nの収集",
     "手動で各サイトを確認\n見落としが発生",
     "AI自動収集・重要度スコアで\n優先情報を即把握"),
    ("SNSの\n活用",
     "個人任せ・体系的な\n把握が困難",
     "地図で一目把握\nリアルタイム更新"),
    ("情報の\n一元管理",
     "情報が分散し\n共有の手間が大きい",
     "ダッシュボードで集約\nチームで即共有"),
    ("ネタ発掘",
     "リサーチャーの\n経験・勘に依存",
     "AIが24時間自動モニタリング\n見落としゼロへ"),
]
for r, (issue, before, after) in enumerate(rows):
    y = 2.0 + r * 1.2
    bg  = WHITE if r % 2 == 0 else LIGHT_BG
    for x, w in zip(col_x, col_w):
        rect(slide, x, y, w, 1.1, fill=bg, line_color=BORDER, line_w=0.5)
    txt(slide, issue, col_x[0] + 0.1, y + 0.15, col_w[0] - 0.2, 0.85,
        size=13, bold=True, color=TSS_GREEN, align=PP_ALIGN.CENTER)
    txt(slide, before, col_x[1] + 0.15, y + 0.1, col_w[1] - 0.3, 0.92,
        size=12, color=GRAY)
    txt(slide, after,  col_x[2] + 0.15, y + 0.1, col_w[2] - 0.3, 0.92,
        size=12, bold=True, color=TSS_GREEN)


# ═══════════════════════════════════════════════════════════
# スライド 8: スケジュール・費用
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
rect(slide, 0, 0, 13.33, 0.06, fill=TSS_RED)
header(slide, "導入スケジュール・費用概算")
page_num(slide, 8)

phases = [
    ("2026年\n4月",  "デモ・ヒアリング",          TSS_GREEN),
    ("2026年\n5月",  "要件確定\n本番環境構築",     TSS_GREEN),
    ("2026年\n6月",  "Phase 1\n本番稼働",          TSS_RED),
    ("2026年\n8月",  "Phase 2\nInstagram・広報資料",TSS_GOLD),
    ("2026年\n度内", "Phase 3\nオープンデータ等",   GRAY),
]
for i, (date, label, col) in enumerate(phases):
    x = 0.5 + i * 2.5
    rect(slide, x, 1.4, 2.1, 0.72, fill=col)
    txt(slide, date, x, 1.4, 2.1, 0.72,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, x, 2.12, 2.1, 1.5, fill=LIGHT_BG, line_color=col, line_w=1.5)
    txt(slide, label, x + 0.1, 2.22, 1.9, 1.3,
        size=12, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(phases) - 1:
        txt(slide, "›", x + 2.1, 1.8, 0.4, 0.5,
            size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# 費用テーブル
rect(slide, 0.4, 4.0, 12.53, 0.5, fill=TSS_GREEN)
txt(slide, "費用概算", 0.55, 4.03, 12.0, 0.45,
    size=14, bold=True, color=WHITE)
costs = [
    ("初期開発・環境構築費", "別途お見積もり"),
    ("月額運用費（サーバー・API利用料）", "別途お見積もり"),
    ("保守・機能追加", "別途ご相談"),
]
for i, (item, val) in enumerate(costs):
    y = 4.6 + i * 0.72
    bg = WHITE if i % 2 == 0 else LIGHT_BG
    rect(slide, 0.4, y, 9.2, 0.62, fill=bg, line_color=BORDER, line_w=0.5)
    rect(slide, 9.6, y, 3.33, 0.62, fill=bg, line_color=BORDER, line_w=0.5)
    txt(slide, item, 0.6, y + 0.1, 8.8, 0.45, size=13, color=DARK)
    txt(slide, val,  9.7, y + 0.1, 3.1, 0.45,
        size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# スライド 9: 次のステップ
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# 背景
rect(slide, 0, 0, 13.33, 7.5, fill=TSS_GREEN)
rect(slide, 0, 0, 13.33, 0.08, fill=TSS_RED)
# ゴールドの装飾ライン
rect(slide, 0, 7.42, 13.33, 0.08, fill=TSS_GOLD)

txt(slide, "次のステップ", 0.6, 0.2, 12, 0.85,
    size=30, bold=True, color=WHITE)
rect(slide, 0.6, 1.05, 3.5, 0.04, fill=TSS_GOLD)

steps = [
    ("01", "プロトタイプデモのご案内",
     "2つのシステムが実際に稼働しています。\n実際の画面をご覧いただきながら、機能・操作感を直接ご確認いただけます。"),
    ("02", "ヒアリング",
     "現在のリサーチフロー・優先課題・運用体制をお聞かせください。\n貴社の業務に最適な構成（分離 or 統合）をご提案します。"),
    ("03", "要件定義・お見積もり",
     "ヒアリング内容をもとに、カスタマイズ範囲と\n費用・スケジュールをご提示します。"),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.5 + i * 1.85
    rect(slide, 0.5, y, 0.9, 1.55, fill=TSS_RED)
    txt(slide, num, 0.5, y + 0.38, 0.9, 0.8,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, 1.6, y + 0.1, 11, 0.55,
        size=19, bold=True, color=WHITE)
    rect(slide, 1.6, y + 0.65, 11.0, 0.03, fill=RGBColor(0x00, 0x8A, 0x66))
    txt(slide, desc, 1.6, y + 0.75, 11.0, 0.75,
        size=13, color=RGBColor(0xA8, 0xD5, 0xC2))


# ─── 保存 ──────────────────────────────────────────────────
out = "/Users/tatsuya1970/projects/tss-ai-sand/needs/proposal.pptx"
prs.save(out)
print(f"保存完了: {out}")
